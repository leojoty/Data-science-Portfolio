###############################################################################
# J-A Legacy PROPRIETARY and Not A Contract Deliverable
# Use and Disclosure Limited to U.S. Employees and only in support of AI programs
# Copyright © J-A Legacy 2019-present. Unpublished work - all rights reserved.
# Third party disclosure requires written approval.
###############################################################################
# -*- coding: utf-8 -*-
"""
Created on Mon Dec  3 12:04:40 2024

__author__ = "Johnathan Svoboda-James"
__copyright__ = "Copyright 2024, J-A Legacy"
__license__ = "J-A Legacy Proprietary"
__maintainer__ = "Johnathan Svoboda-James"
__email__ = "Diversity@johnangie.org
__status__ = "Development"
__version__ = "1.0"
"""

import spacy
from pptx import Presentation

# Load NLP model for Named Entity Recognition
nlp = spacy.load("en_core_web_sm")

# File paths
input_pptx = "C:/Users/gospe/Portfolio/docs/SUMS (1).pptx"
output_pptx = "C:/Users/gospe/Portfolio/docs/Redacted_Project_Proposal.pptx"

# Sensitive entity types to redact
entity_types_to_redact = {"PERSON", "ORG", "DATE", "URL"}

# Helper function to find sensitive terms using NLP
def find_sensitive_terms(text):
    doc = nlp(text)
    sensitive_terms = set()
    for ent in doc.ents:
        if ent.label_ in entity_types_to_redact:
            sensitive_terms.add(ent.text)
    return sensitive_terms

# Open the PowerPoint presentation
presentation = Presentation(input_pptx)

# Iterate through all slides
for slide in presentation.slides:
    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        # Redact text in shapes with text frames
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Extract text and find sensitive terms
                    sensitive_terms = find_sensitive_terms(run.text)
                    
                    # Redact each sensitive term
                    for term in sensitive_terms:
                        run.text = run.text.replace(term, "REDACTED")

        # Redact text in table cells
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    # Extract text and find sensitive terms
                    sensitive_terms = find_sensitive_terms(cell.text)
                    
                    # Redact each sensitive term
                    for term in sensitive_terms:
                        cell.text = cell.text.replace(term, "REDACTED")

# Save the redacted PowerPoint
presentation.save(output_pptx)

print(f"Redacted PowerPoint saved as {output_pptx}")
